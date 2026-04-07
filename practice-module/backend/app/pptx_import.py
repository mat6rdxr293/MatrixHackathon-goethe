from __future__ import annotations

import io
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any
from uuid import uuid4

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageChops


BASE_W = 960
BASE_H = 540
MEDIA_DIR = Path(__file__).parent / "data" / "media"


def _clamp(value: float, min_value: float, max_value: float) -> float:
    return max(min_value, min(max_value, value))


def _resolve_soffice() -> str:
    soffice = os.getenv("SOFFICE_PATH") or ""
    if soffice and not Path(soffice).exists():
        soffice = ""
    if not soffice:
        soffice = shutil.which("soffice") or shutil.which("soffice.exe") or ""
    if not soffice:
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for cand in candidates:
            if Path(cand).exists():
                soffice = cand
                break
    if not soffice:
        raise RuntimeError(
            "LibreOffice не найден. Установите LibreOffice и добавьте soffice в PATH."
        )
    return soffice


def _render_pptx_to_pngs(pptx_bytes: bytes, tmp_path: Path, expected: int | None = None) -> list[Path]:
    soffice = _resolve_soffice()
    pptx_path = tmp_path / "deck.pptx"
    pptx_path.write_bytes(pptx_bytes)

    def _convert(out_dir: Path, filter_name: str) -> list[Path]:
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", filter_name, "--outdir", str(out_dir), str(pptx_path)],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=180,
            )
        except Exception as exc:
            raise RuntimeError(
                f"?? ??????? ?????????????? PPTX ? ??????????? ({filter_name}): {exc}"
            ) from exc
        return sorted(out_dir.glob("*.png"), key=_sort_key)

    images = _convert(tmp_path, "png")
    if expected is not None and len(images) < expected:
        alt_dir = tmp_path / "alt"
        alt_dir.mkdir(parents=True, exist_ok=True)
        alt_images = _convert(alt_dir, "png:impress_png_Export")
        if len(alt_images) > len(images):
            images = alt_images
    if not images:
        raise RuntimeError("LibreOffice ?? ???? ??????? ??????????? ???????.")
    return images


def _save_media(data: bytes, ext: str) -> str:
    MEDIA_DIR.mkdir(parents=True, exist_ok=True)
    safe_ext = (ext or "png").lower().strip(".")
    filename = f"{uuid4().hex}.{safe_ext}"
    path = MEDIA_DIR / filename
    path.write_bytes(data)
    return f"/api/media/{filename}"


def _diff_to_rgba(full_img: Image.Image, without_img: Image.Image, threshold: int = 10, gain: int = 5) -> Image.Image:
    full_rgb = full_img.convert("RGB")
    without_rgb = without_img.convert("RGB")
    diff = ImageChops.difference(full_rgb, without_rgb)
    r, g, b = diff.split()
    alpha = ImageChops.lighter(ImageChops.lighter(r, g), b)

    def _amp(v: int) -> int:
        if v < threshold:
            return 0
        return 255 if v * gain > 255 else v * gain

    alpha = alpha.point(_amp)
    rgba = full_rgb.convert("RGBA")
    rgba.putalpha(alpha)
    return rgba


def _text_from_shape(shape) -> str:
    try:
        if not shape.has_text_frame:
            return ""
        text = shape.text_frame.text or ""
        return text.strip()
    except Exception:
        return ""


def _font_size_from_shape(shape) -> int | None:
    try:
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                size = run.font.size
                if size is not None:
                    return int(_clamp(size.pt, 10, 96))
    except Exception:
        return None
    return None


def _hex_from_color(color) -> str | None:
    try:
        rgb = color.rgb
        if rgb is None:
            return None
        return f"#{rgb}"
    except Exception:
        return None


def _text_style_from_shape(shape) -> tuple[str | None, int | None, str | None, str | None]:
    font_name = None
    font_size = None
    color = None
    align = None
    try:
        if shape.has_text_frame:
            tf = shape.text_frame
            if tf.paragraphs:
                para = tf.paragraphs[0]
                if para.alignment == PP_ALIGN.CENTER:
                    align = "center"
                elif para.alignment == PP_ALIGN.RIGHT:
                    align = "right"
                elif para.alignment == PP_ALIGN.LEFT:
                    align = "left"
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    if not font_name and run.font.name:
                        font_name = run.font.name
                    if font_size is None and run.font.size is not None:
                        font_size = int(_clamp(run.font.size.pt, 10, 96))
                    if color is None and run.font.color is not None:
                        color = _hex_from_color(run.font.color)
                if font_name or font_size or color:
                    break
    except Exception:
        return None, None, None, None
    return font_name, font_size, color, align


def _shape_fill_stroke(shape) -> tuple[str | None, str | None, float | None]:
    fill = None
    stroke = None
    stroke_w = None
    try:
        fill_type = shape.fill.type
        if fill_type in {
            MSO_FILL_TYPE.SOLID,
            MSO_FILL_TYPE.GRADIENT,
            MSO_FILL_TYPE.PICTURE,
            MSO_FILL_TYPE.TEXTURED,
            MSO_FILL_TYPE.PATTERNED,
        }:
            fill = _hex_from_color(shape.fill.fore_color)
    except Exception:
        fill = None
    try:
        line = shape.line
        if line is not None and line.width is not None and line.width.pt and line.width.pt > 0:
            line_fill_type = getattr(line, "fill", None)
            if line_fill_type is None or line_fill_type.type != MSO_FILL_TYPE.BACKGROUND:
                stroke = _hex_from_color(line.color)
            stroke_w = float(line.width.pt)
    except Exception:
        stroke = None
        stroke_w = None
    return fill, stroke, stroke_w


def _shape_image_src(shape) -> tuple[str, str] | None:
    try:
        image = shape.image
    except Exception:
        return None
    try:
        ext = (image.ext or "png").lower()
        src = _save_media(image.blob, ext)
        return src, ext
    except Exception:
        return None


def _group_transform(group) -> tuple[float, float, float, float, float, float]:
    try:
        xfrm = group._element.xfrm  # type: ignore[attr-defined]
        if xfrm is None:
            return 0.0, 0.0, 1.0, 1.0, 0.0, 0.0
        off_x = float(getattr(xfrm.off, "x", 0))
        off_y = float(getattr(xfrm.off, "y", 0))
        ext_cx = float(getattr(xfrm.ext, "cx", 0))
        ext_cy = float(getattr(xfrm.ext, "cy", 0))
        ch_off_x = float(getattr(xfrm.chOff, "x", 0))
        ch_off_y = float(getattr(xfrm.chOff, "y", 0))
        ch_ext_cx = float(getattr(xfrm.chExt, "cx", 0))
        ch_ext_cy = float(getattr(xfrm.chExt, "cy", 0))
        scale_x = ext_cx / ch_ext_cx if ch_ext_cx else 1.0
        scale_y = ext_cy / ch_ext_cy if ch_ext_cy else 1.0
        return off_x, off_y, scale_x, scale_y, ch_off_x, ch_off_y
    except Exception:
        return 0.0, 0.0, 1.0, 1.0, 0.0, 0.0


def _iter_shapes(
    shapes,
    offset_x: float = 0,
    offset_y: float = 0,
    scale_x: float = 1.0,
    scale_y: float = 1.0,
    ch_off_x: float = 0.0,
    ch_off_y: float = 0.0,
):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            g_off_x, g_off_y, g_scale_x, g_scale_y, g_ch_off_x, g_ch_off_y = _group_transform(shape)
            combined_off_x = offset_x + (g_off_x - ch_off_x) * scale_x
            combined_off_y = offset_y + (g_off_y - ch_off_y) * scale_y
            combined_scale_x = scale_x * g_scale_x
            combined_scale_y = scale_y * g_scale_y
            for sub in _iter_shapes(
                shape.shapes,
                combined_off_x,
                combined_off_y,
                combined_scale_x,
                combined_scale_y,
                g_ch_off_x,
                g_ch_off_y,
            ):
                yield sub
        else:
            yield shape, offset_x, offset_y, scale_x, scale_y, ch_off_x, ch_off_y


def import_pptx(pptx_bytes: bytes) -> list[dict[str, Any]]:
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_w = float(prs.slide_width or BASE_W)
    slide_h = float(prs.slide_height or BASE_H)
    scale_x = BASE_W / slide_w
    scale_y = BASE_H / slide_h

    slides: list[dict[str, Any]] = []

    for index, slide in enumerate(prs.slides, start=1):
        elements: list[dict[str, Any]] = []
        title = ""
        background: str | None = None

        for shape, off_x, off_y, grp_scale_x, grp_scale_y, grp_ch_off_x, grp_ch_off_y in _iter_shapes(slide.shapes):
            left = float(getattr(shape, "left", 0))
            top = float(getattr(shape, "top", 0))
            width = float(getattr(shape, "width", 0))
            height = float(getattr(shape, "height", 0))
            world_left = off_x + (left - grp_ch_off_x) * grp_scale_x
            world_top = off_y + (top - grp_ch_off_y) * grp_scale_y
            world_w = width * grp_scale_x
            world_h = height * grp_scale_y
            x = world_left * scale_x
            y = world_top * scale_y
            w = float(world_w) * scale_x
            h = float(world_h) * scale_y
            rotation = 0.0
            try:
                rotation = float(getattr(shape, "rotation", 0.0))
            except Exception:
                rotation = 0.0
            pad = 0
            pad_left = 0
            pad_right = 0
            pad_top = 0
            pad_bottom = 0
            try:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    pad_left = float(getattr(tf, "margin_left", 0)) * scale_x
                    pad_right = float(getattr(tf, "margin_right", 0)) * scale_x
                    pad_top = float(getattr(tf, "margin_top", 0)) * scale_y
                    pad_bottom = float(getattr(tf, "margin_bottom", 0)) * scale_y
                    pad = int(max(0, min(pad_left, pad_right, pad_top, pad_bottom)))
            except Exception:
                pad = 0
                pad_left = 0
                pad_right = 0
                pad_top = 0
                pad_bottom = 0

            img_info = _shape_image_src(shape)
            if img_info:
                src, _ext = img_info
                is_full_bg = (
                    w >= BASE_W * 0.96
                    and h >= BASE_H * 0.96
                    and x <= BASE_W * 0.03
                    and y <= BASE_H * 0.03
                )
                if is_full_bg:
                    background = src
                    continue
                elements.append(
                    {
                        "id": f"img-{index}-{len(elements)+1}",
                        "type": "image",
                        "x": x,
                        "y": y,
                        "w": max(10, w),
                        "h": max(10, h),
                        "src": src,
                        "rotation": rotation,
                    }
                )
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shape_kind = None
                try:
                    if shape.auto_shape_type == MSO_SHAPE.RECTANGLE:
                        shape_kind = "rect"
                    elif shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                        shape_kind = "round"
                    elif shape.auto_shape_type == MSO_SHAPE.OVAL:
                        shape_kind = "ellipse"
                except Exception:
                    shape_kind = None
                if not shape_kind:
                    # Unsupported complex shapes are handled via slide background
                    shape_kind = None
                if not shape_kind:
                    pass
                fill, stroke, stroke_w = _shape_fill_stroke(shape)
                if shape_kind and (fill or stroke):
                    elements.append(
                        {
                            "id": f"shape-{index}-{len(elements)+1}",
                            "type": "shape",
                            "shape": shape_kind,
                            "x": x,
                            "y": y,
                            "w": max(20, w),
                            "h": max(20, h),
                            "fill": fill,
                            "strokeColor": stroke,
                            "strokeWidth": stroke_w,
                            "rotation": rotation,
                        }
                    )

            text = _text_from_shape(shape)
            if text:
                if not title:
                    title = text.splitlines()[0][:60]
                font_name, font_size, font_color, align = _text_style_from_shape(shape)
                elements.append(
                    {
                        "id": f"text-{index}-{len(elements)+1}",
                        "type": "text",
                        "x": x,
                        "y": y,
                        "w": max(40, w),
                        "h": max(30, h),
                        "text": text,
                        "fontSize": font_size or _font_size_from_shape(shape) or 22,
                        "fontFamily": font_name or "Calibri",
                        "color": font_color or "#E7F2FF",
                        "align": align or "left",
                        "padding": pad,
                        "paddingLeft": int(pad_left),
                        "paddingRight": int(pad_right),
                        "paddingTop": int(pad_top),
                        "paddingBottom": int(pad_bottom),
                        "rotation": rotation,
                    }
                )

        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""

        slides.append(
            {
                "id": index,
                "title": title or f"Слайд {index}",
                "content": "",
                "notes": notes,
                "elements": elements,
                "background": background,
            }
        )

    return slides


def _remove_shape_by_id(shapes, target_id: int) -> bool:
    for shape in list(shapes):
        try:
            if shape.shape_id == target_id:
                shape._element.getparent().remove(shape._element)
                return True
        except Exception:
            pass
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if _remove_shape_by_id(shape.shapes, target_id):
                return True
    return False


def _is_complex_shape(shape) -> bool:
    if shape.shape_type in {
        MSO_SHAPE_TYPE.IGX_GRAPHIC,
        MSO_SHAPE_TYPE.DIAGRAM,
        MSO_SHAPE_TYPE.CHART,
        MSO_SHAPE_TYPE.TABLE,
        MSO_SHAPE_TYPE.TEXT_EFFECT,
        MSO_SHAPE_TYPE.FREEFORM,
        MSO_SHAPE_TYPE.LINE,
        MSO_SHAPE_TYPE.CANVAS,
        MSO_SHAPE_TYPE.MEDIA,
        MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
        MSO_SHAPE_TYPE.WEB_VIDEO,
        MSO_SHAPE_TYPE.FORM_CONTROL,
        MSO_SHAPE_TYPE.COMMENT,
        MSO_SHAPE_TYPE.INK,
        MSO_SHAPE_TYPE.INK_COMMENT,
        MSO_SHAPE_TYPE.SCRIPT_ANCHOR,
        MSO_SHAPE_TYPE.CALLOUT,
    }:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            if shape.auto_shape_type in {MSO_SHAPE.RECTANGLE, MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.OVAL}:
                return False
        except Exception:
            return True
        return True
    return False


def import_pptx_stickers(pptx_bytes: bytes) -> list[dict[str, Any]]:
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_w = float(prs.slide_width or BASE_W)
    slide_h = float(prs.slide_height or BASE_H)
    scale_x = BASE_W / slide_w
    scale_y = BASE_H / slide_h

    with tempfile.TemporaryDirectory() as tmp_full:
        tmp_full_path = Path(tmp_full)
        full_paths = _render_pptx_to_pngs(pptx_bytes, tmp_full_path, expected=len(prs.slides))
        full_images: list[Image.Image] = []
        for img_path in full_paths:
            with Image.open(img_path) as img:
                full_images.append(img.convert("RGB").copy())

        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(prs.slides, start=1):
            elements: list[dict[str, Any]] = []
            complex_items: list[dict[str, Any]] = []
            title = ""

            for shape, off_x, off_y, grp_scale_x, grp_scale_y, grp_ch_off_x, grp_ch_off_y in _iter_shapes(
                slide.shapes
            ):
                left = float(getattr(shape, "left", 0))
                top = float(getattr(shape, "top", 0))
                width = float(getattr(shape, "width", 0))
                height = float(getattr(shape, "height", 0))
                world_left = off_x + (left - grp_ch_off_x) * grp_scale_x
                world_top = off_y + (top - grp_ch_off_y) * grp_scale_y
                world_w = width * grp_scale_x
                world_h = height * grp_scale_y
                x = world_left * scale_x
                y = world_top * scale_y
                w = float(world_w) * scale_x
                h = float(world_h) * scale_y

                if w < 1 or h < 1:
                    continue

                rotation = 0.0
                try:
                    rotation = float(getattr(shape, "rotation", 0.0))
                except Exception:
                    rotation = 0.0

                img_info = _shape_image_src(shape)
                if img_info:
                    src, _ext = img_info
                    is_full_bg = (
                        w >= BASE_W * 0.96
                        and h >= BASE_H * 0.96
                        and x <= BASE_W * 0.03
                        and y <= BASE_H * 0.03
                    )
                    if not is_full_bg:
                        elements.append(
                            {
                                "id": f"img-{index}-{len(elements)+1}",
                                "type": "image",
                                "x": x,
                                "y": y,
                                "w": max(10, w),
                                "h": max(10, h),
                                "src": src,
                                "rotation": rotation,
                            }
                        )
                    continue

                if _is_complex_shape(shape):
                    complex_items.append(
                        {
                            "shape_id": shape.shape_id,
                            "x": x,
                            "y": y,
                            "w": w,
                            "h": h,
                            "rotation": rotation,
                            "slide_index": index,
                        }
                    )
                    continue

                # images handled above for any shape that exposes .image

                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    shape_kind = None
                    try:
                        if shape.auto_shape_type == MSO_SHAPE.RECTANGLE:
                            shape_kind = "rect"
                        elif shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                            shape_kind = "round"
                        elif shape.auto_shape_type == MSO_SHAPE.OVAL:
                            shape_kind = "ellipse"
                    except Exception:
                        shape_kind = None
                    if shape_kind:
                        fill, stroke, stroke_w = _shape_fill_stroke(shape)
                        if fill or stroke:
                            elements.append(
                                {
                                    "id": f"shape-{index}-{len(elements)+1}",
                                    "type": "shape",
                                    "shape": shape_kind,
                                    "x": x,
                                    "y": y,
                                    "w": max(20, w),
                                    "h": max(20, h),
                                    "fill": fill,
                                    "strokeColor": stroke,
                                    "strokeWidth": stroke_w,
                                    "rotation": rotation,
                                }
                            )

                text = _text_from_shape(shape)
                if text:
                    if not title:
                        title = text.splitlines()[0][:60]
                    font_name, font_size, font_color, align = _text_style_from_shape(shape)
                    pad_left = pad_right = pad_top = pad_bottom = 0
                    try:
                        if shape.has_text_frame:
                            tf = shape.text_frame
                            pad_left = float(getattr(tf, "margin_left", 0)) * scale_x
                            pad_right = float(getattr(tf, "margin_right", 0)) * scale_x
                            pad_top = float(getattr(tf, "margin_top", 0)) * scale_y
                            pad_bottom = float(getattr(tf, "margin_bottom", 0)) * scale_y
                    except Exception:
                        pass
                    elements.append(
                        {
                            "id": f"text-{index}-{len(elements)+1}",
                            "type": "text",
                            "x": x,
                            "y": y,
                            "w": max(40, w),
                            "h": max(30, h),
                            "text": text,
                            "fontSize": font_size or _font_size_from_shape(shape) or 22,
                            "fontFamily": font_name or "Calibri",
                            "color": font_color or "#E7F2FF",
                            "align": align or "left",
                            "paddingLeft": int(pad_left),
                            "paddingRight": int(pad_right),
                            "paddingTop": int(pad_top),
                            "paddingBottom": int(pad_bottom),
                            "rotation": rotation,
                        }
                    )

            # Generate stickers for complex shapes
            for item in complex_items:
                shape_id = item["shape_id"]
                slide_idx = item["slide_index"]
                # Remove shape and render without it
                prs_without = Presentation(io.BytesIO(pptx_bytes))
                if slide_idx - 1 >= len(prs_without.slides):
                    continue
                slide_without = prs_without.slides[slide_idx - 1]
                removed = _remove_shape_by_id(slide_without.shapes, int(shape_id))
                if not removed:
                    continue
                buf = io.BytesIO()
                prs_without.save(buf)
                with tempfile.TemporaryDirectory() as tmp_without:
                    tmp_path = Path(tmp_without)
                    without_paths = _render_pptx_to_pngs(buf.getvalue(), tmp_path, expected=len(prs_without.slides))
                    if slide_idx - 1 >= len(without_paths):
                        continue
                    with Image.open(without_paths[slide_idx - 1]) as without_img:
                        full_img = full_images[slide_idx - 1]
                        rgba = _diff_to_rgba(full_img, without_img)

                        img_w, img_h = rgba.size
                        sx = img_w / BASE_W
                        sy = img_h / BASE_H
                        pad = 4
                        x_px = int(item["x"] * sx) - pad
                        y_px = int(item["y"] * sy) - pad
                        w_px = int(item["w"] * sx) + pad * 2
                        h_px = int(item["h"] * sy) + pad * 2
                        x_px = max(0, x_px)
                        y_px = max(0, y_px)
                        w_px = max(1, min(img_w - x_px, w_px))
                        h_px = max(1, min(img_h - y_px, h_px))
                        crop = rgba.crop((x_px, y_px, x_px + w_px, y_px + h_px))
                        if crop.getbbox() is None:
                            continue
                        out = io.BytesIO()
                        crop.save(out, format="PNG")
                        src = _save_media(out.getvalue(), "png")
                        elements.append(
                            {
                                "id": f"sticker-{slide_idx}-{shape_id}",
                                "type": "image",
                                "x": item["x"],
                                "y": item["y"],
                                "w": item["w"],
                                "h": item["h"],
                                "src": src,
                                "sticker": True,
                            }
                        )

            notes = ""
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""

            slides.append(
                {
                    "id": index,
                    "title": title or f"Слайд {index}",
                    "content": "",
                    "notes": notes,
                    "elements": elements,
                }
            )

        return slides


def _sort_key(path: Path) -> tuple[int, str]:
    numbers = re.findall(r"\d+", path.stem)
    idx = int(numbers[-1]) if numbers else 0
    return idx, path.name


def import_pptx_full(pptx_bytes: bytes) -> list[dict[str, Any]]:
    # Extract notes using python-pptx for alignment
    prs = Presentation(io.BytesIO(pptx_bytes))
    notes_list: list[str] = []
    for slide in prs.slides:
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        notes_list.append(notes)

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        images = _render_pptx_to_pngs(pptx_bytes, tmp_path, expected=len(prs.slides))

        slides: list[dict[str, Any]] = []
        for i, img_path in enumerate(images, start=1):
            src = _save_media(img_path.read_bytes(), "png")
            notes = notes_list[i - 1] if i - 1 < len(notes_list) else ""
            slides.append(
                {
                    "id": i,
                    "title": f"Слайд {i}",
                    "content": "",
                    "notes": notes,
                    "background": src,
                    "elements": [],
                }
            )

        return slides
