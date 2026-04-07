from __future__ import annotations

import base64
import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt


BASE_W = 960
BASE_H = 540


def _to_emu(value_px: float, total_emu: int, total_px: int) -> int:
    return int((value_px / total_px) * total_emu)


def _data_url_to_bytes(src: str) -> tuple[bytes, str] | None:
    if not src.startswith("data:"):
        return None
    try:
        header, b64 = src.split(",", 1)
        mime = header.split(";")[0].replace("data:", "").strip()
        ext = "png"
        if "jpeg" in mime or "jpg" in mime:
            ext = "jpg"
        elif "png" in mime:
            ext = "png"
        return base64.b64decode(b64), ext
    except Exception:
        return None


def _hex_to_rgb(value: str | None) -> RGBColor | None:
    if not value:
        return None
    v = value.strip()
    if not v.startswith("#"):
        return None
    hex_value = v[1:]
    if len(hex_value) == 3:
        hex_value = "".join([c * 2 for c in hex_value])
    if len(hex_value) != 6:
        return None
    try:
        return RGBColor.from_string(hex_value.upper())
    except Exception:
        return None


def export_pptx(slides: list[dict[str, Any]]) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(int(13.333 * 914400))
    prs.slide_height = Emu(int(7.5 * 914400))

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        background = slide_data.get("background")
        if background:
            img = _data_url_to_bytes(background)
            if img:
                blob, _ = img
                slide.shapes.add_picture(
                    io.BytesIO(blob),
                    Emu(0),
                    Emu(0),
                    width=prs.slide_width,
                    height=prs.slide_height,
                )

        elements = slide_data.get("elements") or []
        if elements:
            for el in elements:
                el_type = el.get("type")
                x = float(el.get("x", 0))
                y = float(el.get("y", 0))
                w = float(el.get("w", 100))
                h = float(el.get("h", 60))
                left = Emu(_to_emu(x, prs.slide_width, BASE_W))
                top = Emu(_to_emu(y, prs.slide_height, BASE_H))
                width = Emu(_to_emu(w, prs.slide_width, BASE_W))
                height = Emu(_to_emu(h, prs.slide_height, BASE_H))

                if el_type == "image":
                    src = el.get("src", "")
                    img = _data_url_to_bytes(src) if isinstance(src, str) else None
                    if img:
                        blob, _ = img
                        slide.shapes.add_picture(
                            io.BytesIO(blob),
                            left,
                            top,
                            width=width,
                            height=height,
                        )
                elif el_type == "text":
                    text = str(el.get("text", ""))
                    font_size = int(el.get("fontSize") or 22)
                    font_family = el.get("fontFamily")
                    color = _hex_to_rgb(el.get("color"))
                    align = str(el.get("align") or "left").lower()
                    box = slide.shapes.add_textbox(left, top, width, height)
                    tf = box.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    if align == "center":
                        p.alignment = PP_ALIGN.CENTER
                    elif align == "right":
                        p.alignment = PP_ALIGN.RIGHT
                    else:
                        p.alignment = PP_ALIGN.LEFT
                    run = p.add_run()
                    run.text = text
                    run.font.size = Pt(font_size)
                    if font_family:
                        run.font.name = str(font_family)
                    if color:
                        run.font.color.rgb = color
                elif el_type == "shape":
                    shape_kind = str(el.get("shape") or "rect")
                    shape_type = {
                        "rect": MSO_SHAPE.RECTANGLE,
                        "round": MSO_SHAPE.ROUNDED_RECTANGLE,
                        "ellipse": MSO_SHAPE.OVAL,
                    }.get(shape_kind, MSO_SHAPE.RECTANGLE)
                    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
                    fill_color = _hex_to_rgb(el.get("fill"))
                    if fill_color:
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = fill_color
                    stroke_color = _hex_to_rgb(el.get("strokeColor"))
                    stroke_width = float(el.get("strokeWidth") or 0)
                    if stroke_color and stroke_width > 0:
                        shp.line.color.rgb = stroke_color
                        shp.line.width = Pt(stroke_width)
                    elif stroke_width <= 0:
                        shp.line.fill.background()
        else:
            content = slide_data.get("content") or ""
            if content:
                left = Emu(_to_emu(60, prs.slide_width, BASE_W))
                top = Emu(_to_emu(60, prs.slide_height, BASE_H))
                width = Emu(_to_emu(BASE_W - 120, prs.slide_width, BASE_W))
                height = Emu(_to_emu(BASE_H - 120, prs.slide_height, BASE_H))
                box = slide.shapes.add_textbox(left, top, width, height)
                tf = box.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(content)
                run.font.size = Pt(28)

        notes = slide_data.get("notes") or ""
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = str(notes)
            except Exception:
                pass

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
